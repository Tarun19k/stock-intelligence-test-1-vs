"""
GSI Dashboard — Community Launch Pitch Deck
Run: python3 docs/build_pitch_deck.py
Output: docs/GSI_Community_Launch_Deck.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Palette ──────────────────────────────────────────────────────────────────
BG        = RGBColor(0x0d, 0x11, 0x17)   # near-black (matches app)
BG_CARD   = RGBColor(0x16, 0x1b, 0x27)   # card surface
ACCENT    = RGBColor(0x4f, 0xc3, 0xf7)   # electric blue
WHITE     = RGBColor(0xff, 0xff, 0xff)
MUTED     = RGBColor(0x7e, 0x8b, 0xa0)
GREEN     = RGBColor(0x2e, 0x7d, 0x32)
GREEN_TXT = RGBColor(0x66, 0xbb, 0x6a)
AMBER     = RGBColor(0xff, 0x8f, 0x00)
RED_TXT   = RGBColor(0xef, 0x53, 0x50)

W, H = Inches(13.33), Inches(7.5)   # 16:9 widescreen

# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)

def bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, x, y, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape

def txt(slide, text, x, y, w, h, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, font="Inter"):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    run.font.name    = font
    return txb

def label(slide, text, x, y, w=2.5, color=ACCENT, size=9):
    """Small all-caps label / eyebrow"""
    txt(slide, text.upper(), x, y, w, 0.3, size=size, bold=True,
        color=color, font="JetBrains Mono")

def divider(slide, y, x=0.6, w=12.1, color=ACCENT, thick=0.03):
    bar = slide.shapes.add_shape(1,
        Inches(x), Inches(y), Inches(w), Inches(thick))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = color

def pill(slide, text, x, y, w, h, bg_color, txt_color=WHITE, size=11):
    box(slide, x, y, w, h, bg_color)
    txt(slide, text, x + 0.05, y + 0.02, w - 0.1, h - 0.04,
        size=size, bold=True, color=txt_color, align=PP_ALIGN.CENTER)

def slide_number(slide, n, total=9):
    txt(slide, f"{n} / {total}", 12.5, 7.1, 0.7, 0.3,
        size=8, color=MUTED, align=PP_ALIGN.RIGHT)

# ── Slides ────────────────────────────────────────────────────────────────────

def slide_01_title(prs):
    """Title — hero"""
    s = blank_slide(prs)
    bg(s)

    # Accent left bar
    box(s, 0, 0, 0.08, 7.5, ACCENT)

    # Eyebrow
    label(s, "Community Launch  •  Free  •  No Login", 0.5, 1.1, 10)

    # Headline
    txt(s, "Multi-market stock signals.\nTwo frameworks. One verdict.",
        0.5, 1.6, 9, 2.2, size=44, bold=True, color=WHITE)

    # Accent underline
    divider(s, 3.95, x=0.5, w=5.5, color=ACCENT, thick=0.04)

    # Sub
    txt(s, "Weinstein Stage + Elder Triple Screen → BUY / WATCH / AVOID\nacross 9 markets, 556 tickers — free, open-source, no account.",
        0.5, 4.1, 8, 1.0, size=16, color=MUTED)

    # Stats row
    stats = [("9", "Markets"), ("556", "Tickers"), ("38", "Groups"), ("Free", "Always")]
    for i, (val, lbl) in enumerate(stats):
        xpos = 0.5 + i * 2.2
        txt(s, val,  xpos, 5.3, 2.0, 0.7, size=36, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        txt(s, lbl,  xpos, 5.95, 2.0, 0.3, size=11, color=MUTED,  align=PP_ALIGN.CENTER)

    # Right: app screenshot placeholder
    box(s, 9.5, 1.0, 3.5, 5.5, BG_CARD)
    txt(s, "[ App Screenshot ]", 9.5, 3.8, 3.5, 0.5, size=11,
        color=MUTED, align=PP_ALIGN.CENTER, italic=True)

    slide_number(s, 1)

def slide_02_problem(prs):
    """Problem"""
    s = blank_slide(prs)
    bg(s)
    box(s, 0, 0, 0.08, 7.5, RED_TXT)

    label(s, "The Problem", 0.5, 0.4)
    txt(s, "Retail investors face a bad trade-off.",
        0.5, 0.85, 12, 0.7, size=34, bold=True, color=WHITE)
    divider(s, 1.75, x=0.5, w=12.1, color=RED_TXT)

    problems = [
        ("💸", "TradingView Pro", "₹1,200/month for Weinstein Stage + proper indicators.\nMost retail investors hit the free tier ceiling."),
        ("🇺🇸", "Yahoo Finance / Google", "US-centric. No signal aggregation.\nNo Weinstein. No Elder. No unified verdict."),
        ("🇮🇳", "Screener.in / Tickertape", "India-only. Fundamental-heavy.\nNo technical signal framework."),
    ]
    for i, (icon, heading, body) in enumerate(problems):
        xpos = 0.5 + i * 4.2
        box(s, xpos, 2.1, 3.9, 3.8, BG_CARD)
        txt(s, icon,    xpos + 0.2, 2.3,  0.6, 0.6, size=28)
        txt(s, heading, xpos + 0.2, 3.0,  3.5, 0.5, size=15, bold=True, color=WHITE)
        txt(s, body,    xpos + 0.2, 3.55, 3.5, 1.8, size=12, color=MUTED)

    txt(s, "The gap: institutional-grade technical analysis, accessible to everyone, across global markets.",
        0.5, 6.3, 12.1, 0.5, size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    slide_number(s, 2)

def slide_03_solution(prs):
    """Solution"""
    s = blank_slide(prs)
    bg(s)
    box(s, 0, 0, 0.08, 7.5, GREEN)

    label(s, "The Solution", 0.5, 0.4)
    txt(s, "GSI Dashboard — free, open-source, no login.",
        0.5, 0.85, 12, 0.7, size=34, bold=True, color=WHITE)
    divider(s, 1.75, x=0.5, w=12.1, color=GREEN_TXT)

    features = [
        ("📊", "4-Tab Stock Dashboard",     "Charts · Forecast · Compare · Insights\nFor every ticker in your watchlist."),
        ("🌍", "Global Intelligence",        "Geopolitical events → sector impact chains.\nMacro context with one click."),
        ("📅", "Week Summary",               "Market + group overview every week.\nTop movers. Signal heatmap."),
        ("⚖️",  "Portfolio Allocator",        "Mean-CVaR optimisation. Stress regime\ndetection. No Bloomberg required."),
        ("📈", "Live Ticker Bar",            "Real-time price movement across\nall markets. Refreshes every 5s."),
        ("🔔", "Forecast Tracker",           "Log your directional forecasts.\nTrack accuracy over time."),
    ]
    for i, (icon, heading, body) in enumerate(features):
        col = i % 3
        row = i // 3
        xpos = 0.5 + col * 4.2
        ypos = 2.1 + row * 2.35
        box(s, xpos, ypos, 3.9, 2.1, BG_CARD)
        txt(s, icon,    xpos + 0.15, ypos + 0.15, 0.5, 0.5, size=20)
        txt(s, heading, xpos + 0.7,  ypos + 0.15, 3.1, 0.45, size=13, bold=True, color=WHITE)
        txt(s, body,    xpos + 0.15, ypos + 0.65, 3.6, 1.2,  size=11, color=MUTED)

    slide_number(s, 3)

def slide_04_markets(prs):
    """Markets covered"""
    s = blank_slide(prs)
    bg(s)
    box(s, 0, 0, 0.08, 7.5, ACCENT)

    label(s, "Market Coverage", 0.5, 0.4)
    txt(s, "9 markets. One dashboard. Same signal engine.",
        0.5, 0.85, 12, 0.7, size=34, bold=True, color=WHITE)
    divider(s, 1.75, x=0.5, w=12.1, color=ACCENT)

    markets = [
        ("🇮🇳", "India",        "Nifty 50 · Next 50 · Midcap · Smallcap\nBank Nifty · IT · Pharma · FMCG"),
        ("🇺🇸", "USA",          "S&P 500 · Nasdaq · Dow\nTech · Financials · Energy"),
        ("🇪🇺", "Europe",       "DAX · CAC 40 · FTSE 100\nStoxx 600 components"),
        ("🇨🇳", "China",        "Shanghai · Shenzhen\nHang Seng · H-Shares"),
        ("🌏", "Asia-Pacific", "Nikkei · ASX 200\nSingapore · Korea"),
        ("🌍", "West Asia",    "Saudi Aramco · ADNOC\nQatar · UAE markets"),
        ("🥇", "Commodities",  "Gold · Silver · Crude Oil\nNatural Gas · Copper"),
        ("💱", "Forex",        "USD/INR · EUR/USD\nGBP/USD · JPY pairs"),
        ("📊", "Rates & Bonds","US 10Y · India 10Y\nVIX · Fear & Greed"),
    ]
    for i, (flag, name, tickers) in enumerate(markets):
        col = i % 3
        row = i // 3
        xpos = 0.5 + col * 4.2
        ypos = 2.15 + row * 1.65
        box(s, xpos, ypos, 3.9, 1.5, BG_CARD)
        txt(s, f"{flag}  {name}", xpos + 0.15, ypos + 0.12, 3.6, 0.4,
            size=14, bold=True, color=WHITE)
        txt(s, tickers, xpos + 0.15, ypos + 0.55, 3.6, 0.85, size=10, color=MUTED)

    slide_number(s, 4)

def slide_05_signal_engine(prs):
    """Signal engine"""
    s = blank_slide(prs)
    bg(s)
    box(s, 0, 0, 0.08, 7.5, ACCENT)

    label(s, "Signal Engine", 0.5, 0.4)
    txt(s, "Not one method. Three — arbitrated.",
        0.5, 0.85, 12, 0.7, size=34, bold=True, color=WHITE)
    divider(s, 1.75, x=0.5, w=12.1, color=ACCENT)

    # Three method cards
    methods = [
        ("WEINSTEIN",  ACCENT,    "Stage 1–4 cycle position.\nLong-term trend.\nTakes precedence in conflict."),
        ("ELDER",      AMBER,     "Triple Screen — weekly trend,\ndaily momentum,\nintraday entry timing."),
        ("COMPOSITE",  MUTED,     "RSI · MACD · Volume · ATR\n+ 8 more indicators.\nScored 0–100."),
    ]
    for i, (name, color, desc) in enumerate(methods):
        xpos = 0.5 + i * 4.1
        box(s, xpos, 2.1, 3.8, 3.0, BG_CARD)
        pill(s, name, xpos + 0.15, 2.25, 1.5, 0.38, color, WHITE, size=10)
        txt(s, desc, xpos + 0.15, 2.8, 3.5, 2.1, size=12, color=MUTED)

    # Arrow → verdict
    txt(s, "↓", 6.3, 5.2, 0.8, 0.5, size=24, color=ACCENT, align=PP_ALIGN.CENTER)

    # Verdict
    box(s, 3.5, 5.75, 6.3, 1.1, GREEN)
    txt(s, "▲  UNIFIED VERDICT  →  BUY", 3.5, 5.75, 6.3, 1.1,
        size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txt(s, "Documented override hierarchy: Weinstein > Elder in conflict. Veto disclosed in UI.",
        0.5, 7.0, 12.1, 0.35, size=10, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

    slide_number(s, 5)

def slide_06_comparison(prs):
    """Free vs competitors"""
    s = blank_slide(prs)
    bg(s)
    box(s, 0, 0, 0.08, 7.5, ACCENT)

    label(s, "Competitive Landscape", 0.5, 0.4)
    txt(s, "Free doesn't mean stripped down.",
        0.5, 0.85, 12, 0.7, size=34, bold=True, color=WHITE)
    divider(s, 1.75, x=0.5, w=12.1, color=ACCENT)

    # Table header
    cols     = ["Feature", "GSI Dashboard", "TradingView Free", "Yahoo Finance", "Screener.in"]
    col_w    = [3.5, 2.2, 2.2, 2.2, 2.2]
    col_x    = [0.5, 4.1, 6.3, 8.5, 10.7]
    row_h    = 0.52
    y_start  = 2.15

    # Header row
    for j, (heading, cw, cx) in enumerate(zip(cols, col_w, col_x)):
        hdr_color = ACCENT if j == 1 else BG_CARD
        box(s, cx, y_start, cw, row_h, hdr_color)
        txt(s, heading, cx + 0.1, y_start + 0.08, cw - 0.2, row_h - 0.1,
            size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

    rows = [
        ("Weinstein Stage",         "✓", "Pro only",  "✗",  "✗"),
        ("Elder Triple Screen",     "✓", "Pro only",  "✗",  "✗"),
        ("Unified BUY/WATCH/AVOID", "✓", "✗",         "✗",  "✗"),
        ("India markets",           "✓", "✓",          "Partial", "✓"),
        ("9 global markets",        "✓", "✓ (limited)", "✗",  "✗"),
        ("Portfolio optimiser",     "✓", "Pro only",  "✗",  "✗"),
        ("No login required",       "✓", "✓",          "✓",  "✓"),
        ("Cost",                    "Free", "₹0–1,200/mo", "Free", "Free"),
    ]

    for r, (feat, *vals) in enumerate(rows):
        ypos = y_start + (r + 1) * row_h
        row_bg = BG if r % 2 == 0 else BG_CARD
        for j, (val, cw, cx) in enumerate(zip([feat] + vals, col_w, col_x)):
            cell_bg = RGBColor(0x0d, 0x28, 0x1a) if j == 1 else row_bg
            box(s, cx, ypos, cw, row_h, cell_bg)
            val_color = GREEN_TXT if val == "✓" or val == "Free" else (RED_TXT if val == "✗" else WHITE)
            txt(s, val, cx + 0.1, ypos + 0.08, cw - 0.2, row_h - 0.1,
                size=11, color=val_color,
                align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

    slide_number(s, 6)

def slide_07_personas(prs):
    """Who it's for"""
    s = blank_slide(prs)
    bg(s)
    box(s, 0, 0, 0.08, 7.5, ACCENT)

    label(s, "Who It's For", 0.5, 0.4)
    txt(s, "Three users. One tool.",
        0.5, 0.85, 12, 0.7, size=34, bold=True, color=WHITE)
    divider(s, 1.75, x=0.5, w=12.1, color=ACCENT)

    personas = [
        ("Ritu", "Indian Retail Investor",
         "Mumbai · Zerodha",
         "NSE/BSE technical signals without paying ₹1,200/month for TradingView Pro.",
         '"I can see the Weinstein stage for my\nwatchlist in one view."'),
        ("Alex", "US Retail Trader",
         "Austin, TX · Robinhood",
         "Multi-market breadth — what's moving globally before the US market opens.",
         '"I can see India/China/Europe signals\nbefore US open — no other free tool does this."'),
        ("Priya", "Finance Student",
         "Bangalore · University",
         "Apply Weinstein Stage and Elder Triple Screen to 556 real tickers for assignments.",
         '"I understand what Stage 2 looks like\non a real chart."'),
    ]
    for i, (name, role, loc, pain, quote) in enumerate(personas):
        xpos = 0.5 + i * 4.2
        box(s, xpos, 2.1, 3.9, 4.8, BG_CARD)
        txt(s, name,  xpos + 0.2, 2.25, 3.5, 0.5, size=18, bold=True, color=ACCENT)
        txt(s, role,  xpos + 0.2, 2.75, 3.5, 0.35, size=11, bold=True, color=WHITE)
        txt(s, loc,   xpos + 0.2, 3.1,  3.5, 0.3,  size=10, color=MUTED)
        divider(s, 3.5, x=xpos + 0.2, w=3.5, color=BG, thick=0.01)
        txt(s, pain,  xpos + 0.2, 3.55, 3.5, 1.0, size=11, color=MUTED)
        box(s, xpos + 0.2, 4.65, 3.5, 1.9, BG)
        txt(s, quote, xpos + 0.3, 4.75, 3.3, 1.7, size=11, italic=True, color=WHITE)

    slide_number(s, 7)

def slide_08_cta(prs):
    """CTA"""
    s = blank_slide(prs)
    bg(s)
    box(s, 0, 0, 0.08, 7.5, GREEN)

    label(s, "Try It Now", 4.5, 1.5, w=4, color=GREEN_TXT)
    txt(s, "Start analysing.\nFree. Right now.",
        1.5, 2.0, 10, 2.0, size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txt(s, "No login. No signup. No subscription. Just open the dashboard.",
        1.5, 4.1, 10, 0.5, size=16, color=MUTED, align=PP_ALIGN.CENTER)

    # CTA buttons
    box(s, 3.2, 4.85, 3.0, 0.7, ACCENT)
    txt(s, "Open Dashboard →", 3.2, 4.85, 3.0, 0.7,
        size=16, bold=True, color=BG, align=PP_ALIGN.CENTER)

    box(s, 6.6, 4.85, 3.0, 0.7, BG_CARD)
    txt(s, "View on GitHub →", 6.6, 4.85, 3.0, 0.7,
        size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txt(s, "Hosted on Streamlit Community Cloud  ·  Free tier  ·  May take 10–15s on cold start",
        1.5, 5.9, 10, 0.35, size=10, color=MUTED, align=PP_ALIGN.CENTER)

    slide_number(s, 8)

def slide_09_legal(prs):
    """Legal / disclaimer — required on all decks"""
    s = blank_slide(prs)
    bg(s)
    box(s, 0, 0, 0.08, 7.5, AMBER)

    label(s, "Important Disclaimer", 0.5, 0.6, color=AMBER)
    divider(s, 1.3, x=0.5, w=12.1, color=AMBER, thick=0.03)

    disclaimer = (
        "This tool is for educational and research purposes only. It is NOT a SEBI-registered investment adviser, "
        "research analyst, or financial planner. Nothing in this dashboard constitutes investment advice, a buy/sell "
        "recommendation, or a solicitation to trade any security.\n\n"
        "All signals are algorithmically generated from quantitative calculations and should be used only as one "
        "input among many in your own research process.\n\n"
        "Data is sourced via yfinance (Yahoo Finance) and carries a 15–20 minute delay for most markets. "
        "Data accuracy is not guaranteed. Verify independently before making any investment decision.\n\n"
        "Use of this tool is subject to Yahoo Finance's Terms of Service.\n\n"
        "NOT SEBI-registered  ·  NOT SEC-registered  ·  NOT FCA-registered  ·  NOT investment advice"
    )
    txt(s, disclaimer, 0.5, 1.5, 12.2, 5.5, size=13, color=MUTED)

    slide_number(s, 9)


# ── Build ─────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    prs = new_prs()
    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_solution(prs)
    slide_04_markets(prs)
    slide_05_signal_engine(prs)
    slide_06_comparison(prs)
    slide_07_personas(prs)
    slide_08_cta(prs)
    slide_09_legal(prs)

    out = os.path.join(os.path.dirname(__file__), "GSI_Community_Launch_Deck.pptx")
    prs.save(out)
    print(f"✅  Saved: {out}")
    print(f"   9 slides  ·  16:9 widescreen  ·  dark theme")
